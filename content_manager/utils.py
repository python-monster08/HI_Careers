import io
from PyPDF2 import PdfFileReader
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file):
    reader = PdfFileReader(file)
    text = ''
    for page in reader.pages:
        text += page.extractText()
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text
